import os
import io
import asyncio
import tempfile
import subprocess
import random
from typing import Dict, Any, List, Optional, Literal
from pathlib import Path

from .image_generator import generate_slide_image, get_placeholder_image

# Define types for image placement
ImagePlacementType = Literal['background', 'side', 'none']

async def create_presentation(
    content_data: Dict[str, Any], 
    template_style: str = "corporate", 
    output_path: Optional[str] = None, 
    use_background_images: bool = True,
    image_placement: ImagePlacementType = "background",
    style_properties: Optional[Dict[str, Any]] = None
) -> str:
    """
    Create a PowerPoint presentation from the generated content.
    
    Args:
        content_data: The presentation content data
        template_style: The template style to use
        output_path: Path to save the presentation (optional)
        use_background_images: Whether to use images as slide backgrounds (legacy option)
        image_placement: How to place images in slides ('background', 'side', or 'none')
        style_properties: Dictionary of custom style properties to apply to the presentation
        
    Returns:
        Path to the created presentation file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        # For backward compatibility
        if use_background_images and image_placement == "side":
            image_placement = "background"
        
        # Get style properties from content_data if not provided directly
        if style_properties is None and "style" in content_data:
            style_properties = content_data["style"]
        
        # Create base templates directory if it doesn't exist
        base_template_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates")
        os.makedirs(base_template_dir, exist_ok=True)
        
        # Check if the template exists, if not create a basic one
        template_path = os.path.join(base_template_dir, f"{template_style}.pptx")
        if not os.path.exists(template_path):
            await create_base_template(template_style, template_path, style_properties)
        
        # Load the template
        prs = Presentation(template_path)
        
        # Apply any custom style properties to the presentation
        if style_properties:
            try:
                await apply_custom_style_to_presentation(prs, style_properties)
                print(f"Applied custom style properties to presentation")
            except Exception as e:
                print(f"Error applying custom style: {str(e)}")
        
        # Get slide layouts
        title_layout = prs.slide_layouts[0]  # Title slide
        content_layout = prs.slide_layouts[1]  # Content slide with title and content
        
        # Track created slides
        slides = []
        
        # Create slides based on the content data
        if "slides" in content_data:
            for slide_data in content_data["slides"]:
                # Make a copy of slide_data to avoid modifying the original
                processed_slide_data = slide_data.copy()
                
                # Extract image prompt but don't include it in the final slide
                image_prompt = processed_slide_data.pop("image_prompt", None)
                
                # Skip image processing if image_placement is 'none'
                include_image = image_placement != "none" and image_prompt
                
                # Determine if we should use background style
                use_as_background = image_placement == "background"
                
                # For side placement, randomly choose left or right
                image_position = "center"
                if image_placement == "side":
                    image_position = random.choice(["left", "right"])
                
                if processed_slide_data.get("type") == "title":
                    # Create title slide
                    slide = prs.slides.add_slide(title_layout)
                    slides.append(slide)
                    
                    if "title" in processed_slide_data:
                        slide.shapes.title.text = processed_slide_data["title"]
                    
                    # Handle subtitle if present
                    subtitle = processed_slide_data.get("subtitle") or processed_slide_data.get("content")
                    if subtitle and len(slide.placeholders) > 1:
                        slide.placeholders[1].text = subtitle
                    
                    # Generate and add image if image prompt is specified
                    if include_image:
                        print(f"Generating image for title slide with prompt: {image_prompt}")
                        image_url = await generate_slide_image(image_prompt, template_style)
                        if image_url:
                            await add_image_to_slide(
                                slide, 
                                image_url, 
                                position="center" if use_as_background else image_position, 
                                as_background=use_as_background
                            )
                        else:
                            # Use placeholder image if generation fails
                            placeholder_url = get_placeholder_image(template_style)
                            if os.path.exists(placeholder_url):
                                await add_image_to_slide(
                                    slide, 
                                    placeholder_url, 
                                    position="center" if use_as_background else image_position, 
                                    is_local=True, 
                                    as_background=use_as_background
                                )
                
                elif processed_slide_data.get("type") in ["content", "bullets"]:
                    # Create content slide
                    slide = prs.slides.add_slide(content_layout)
                    slides.append(slide)
                    
                    if "title" in processed_slide_data:
                        slide.shapes.title.text = processed_slide_data["title"]
                    
                    # Add bullet points if present
                    bullets = processed_slide_data.get("bullets", [])
                    if bullets and len(slide.placeholders) > 1:
                        text_frame = slide.placeholders[1].text_frame
                        text_frame.clear()
                        
                        for bullet in bullets:
                            p = text_frame.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    
                    # Add general content if present and no bullets
                    content = processed_slide_data.get("content")
                    if content and not bullets and len(slide.placeholders) > 1:
                        slide.placeholders[1].text = content
                    
                    # Generate and add image if image prompt is specified
                    if include_image:
                        print(f"Generating image for content slide '{processed_slide_data.get('title', '')}' with prompt: {image_prompt}")
                        image_url = await generate_slide_image(image_prompt, template_style)
                        if image_url:
                            if use_as_background:
                                await add_image_to_slide(slide, image_url, as_background=True)
                            else:
                                await add_image_to_slide(slide, image_url, position=image_position)
                
                elif processed_slide_data.get("type") == "image":
                    # Create image-focused slide
                    slide = prs.slides.add_slide(content_layout)
                    slides.append(slide)
                    
                    if "title" in processed_slide_data:
                        slide.shapes.title.text = processed_slide_data["title"]
                    
                    # Add caption/content if present
                    content = processed_slide_data.get("content")
                    if content and len(slide.placeholders) > 1:
                        slide.placeholders[1].text = content
                    
                    # Generate and add image if image prompt is specified
                    if include_image:
                        print(f"Generating image for image slide with prompt: {image_prompt}")
                        image_url = await generate_slide_image(image_prompt, template_style)
                        if image_url:
                            await add_image_to_slide(
                                slide, 
                                image_url, 
                                position="center" if use_as_background else image_position, 
                                size="large" if not use_as_background else "medium", 
                                as_background=use_as_background
                            )
                
                elif processed_slide_data.get("type") == "quote":
                    # Create quote slide
                    slide = prs.slides.add_slide(content_layout)
                    slides.append(slide)
                    
                    if "title" in processed_slide_data:
                        slide.shapes.title.text = processed_slide_data["title"]
                    
                    # Add quote text
                    if "quote" in processed_slide_data and len(slide.placeholders) > 1:
                        text_frame = slide.placeholders[1].text_frame
                        text_frame.clear()
                        
                        p = text_frame.add_paragraph()
                        p.text = f'"{processed_slide_data["quote"]}"'
                        
                        # Add author if present
                        if "author" in processed_slide_data:
                            p = text_frame.add_paragraph()
                            p.text = f"— {processed_slide_data['author']}"
                            p.level = 1
                    
                    # Generate and add image if image prompt is specified
                    if include_image:
                        print(f"Generating image for quote slide with prompt: {image_prompt}")
                        image_url = await generate_slide_image(image_prompt, template_style)
                        if image_url:
                            await add_image_to_slide(
                                slide, 
                                image_url, 
                                position=image_position, 
                                as_background=use_as_background
                            )
                
                elif processed_slide_data.get("type") == "two-column":
                    # Create two-column slide
                    slide = prs.slides.add_slide(content_layout)
                    slides.append(slide)
                    
                    if "title" in processed_slide_data:
                        slide.shapes.title.text = processed_slide_data["title"]
                    
                    # Add left column content
                    left_content = processed_slide_data.get("left_content", "")
                    right_content = processed_slide_data.get("right_content", "")
                    
                    # Use text placeholder for combined content
                    if len(slide.placeholders) > 1 and (left_content or right_content):
                        text_frame = slide.placeholders[1].text_frame
                        text_frame.clear()
                        
                        if left_content:
                            p_left = text_frame.add_paragraph()
                            p_left.text = "Left Column:"
                            p_left.level = 0
                            
                            p_left_content = text_frame.add_paragraph()
                            p_left_content.text = left_content
                            p_left_content.level = 1
                        
                        if right_content:
                            p_right = text_frame.add_paragraph()
                            p_right.text = "Right Column:"
                            p_right.level = 0
                            
                            p_right_content = text_frame.add_paragraph()
                            p_right_content.text = right_content
                            p_right_content.level = 1
                    
                    # Generate and add image if image prompt is specified
                    if include_image:
                        print(f"Generating image for two-column slide with prompt: {image_prompt}")
                        image_url = await generate_slide_image(image_prompt, template_style)
                        if image_url:
                            # For two-column slides, prefer bottom or top position when not background
                            side_position = "bottom" if not use_as_background else image_position
                            await add_image_to_slide(
                                slide, 
                                image_url, 
                                position=side_position, 
                                as_background=use_as_background
                            )
        
        # Define output path if not provided
        if not output_path:
            output_path = os.path.join(tempfile.gettempdir(), f"presentation_{template_style}.pptx")
        
        # Save the presentation
        prs.save(output_path)
        return output_path
        
    except ImportError:
        print("python-pptx library not installed. Install it with: pip install python-pptx")
        raise
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")
        raise

async def add_image_to_slide(
    slide, 
    image_url: str, 
    position: str = "center",  # center, left, right, top, bottom
    size: str = "medium",      # small, medium, large
    is_local: bool = False,
    as_background: bool = False
):
    """
    Add an image to a slide at the specified position.
    
    Args:
        slide: The slide to add the image to
        image_url: URL or local path of the image
        position: Where to position the image (center, left, right, top, bottom)
        size: Size of the image (small, medium, large)
        is_local: Whether the image is a local file
        as_background: Whether to use the image as a background
    """
    try:
        from pptx.util import Inches
        import requests
        from PIL import Image
        
        # Get slide dimensions for reference
        slide_width = slide.slide_width.inches
        slide_height = slide.slide_height.inches
        
        # Define image size based on input
        size_map = {
            "small": (2, 2),  # 2x2 inches
            "medium": (4, 3),  # 4x3 inches
            "large": (6, 4.5),  # 6x4.5 inches
        }
        
        # Default to medium size
        img_width, img_height = size_map.get(size, size_map["medium"])
        
        # Get the image data
        if is_local:
            with open(image_url, "rb") as f:
                img_data = f.read()
        else:
            response = requests.get(image_url)
            img_data = response.content
        
        # If using as background, fill the entire slide
        if as_background:
            # Create temporary file to store image
            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as img_file:
                img_file.write(img_data)
                img_file_path = img_file.name
            
            # Process the image to create a semi-transparent overlay
            # This helps to ensure text is readable over the background
            try:
                img = Image.open(img_file_path)
                
                # Create a semi-transparent white overlay to improve text readability
                overlay = Image.new("RGBA", img.size, (255, 255, 255, 80))  # 80 is the alpha (transparency)
                
                # If the original image has an alpha channel, convert overlay accordingly
                if img.mode == "RGBA":
                    img = Image.alpha_composite(img.convert("RGBA"), overlay)
                else:
                    # Convert to RGBA first
                    img = img.convert("RGBA")
                    overlay_img = Image.alpha_composite(img, overlay)
                    img = overlay_img
                
                # Save the processed image
                img.convert("RGB").save(img_file_path)
                
                # Add as background
                slide.shapes.add_picture(img_file_path, 0, 0, width=Inches(slide_width), height=Inches(slide_height))
                
                # Clean up
                os.unlink(img_file_path)
            
            except Exception as e:
                print(f"Error processing background image: {str(e)}")
                # Fallback to direct addition without overlay
                slide.shapes.add_picture(img_file_path, 0, 0, width=Inches(slide_width), height=Inches(slide_height))
                os.unlink(img_file_path)
        
        else:
            # Not a background, position according to parameters
            # Calculate position based on input
            if position == "center":
                left = (slide_width - img_width) / 2
                top = (slide_height - img_height) / 2
            elif position == "left":
                left = 0.5  # 0.5 inch margin
                top = (slide_height - img_height) / 2
                # For side images, adjust width to be narrower
                img_width = min(img_width, slide_width / 2.5)
                img_height = img_width * 0.75  # Maintain aspect ratio
            elif position == "right":
                left = slide_width - img_width - 0.5  # 0.5 inch margin
                top = (slide_height - img_height) / 2
                # For side images, adjust width to be narrower
                img_width = min(img_width, slide_width / 2.5)
                img_height = img_width * 0.75  # Maintain aspect ratio
            elif position == "top":
                left = (slide_width - img_width) / 2
                top = 0.5  # 0.5 inch margin
            elif position == "bottom":
                left = (slide_width - img_width) / 2
                top = slide_height - img_height - 0.5  # 0.5 inch margin
            else:
                # Default to center if unknown position
                left = (slide_width - img_width) / 2
                top = (slide_height - img_height) / 2
            
            # Create in-memory stream for image
            image_stream = io.BytesIO(img_data)
            
            # Add image to slide at calculated position
            slide.shapes.add_picture(image_stream, Inches(left), Inches(top), width=Inches(img_width), height=Inches(img_height))
    
    except Exception as e:
        print(f"Error adding image to slide: {str(e)}")
        # Continue without the image

async def apply_custom_style_to_presentation(prs, style_properties):
    """Apply custom style properties to the presentation."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    # Helper function to convert hex color to RGB
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    # Apply style to all slide masters
    for master in prs.slide_masters:
        # Process all slide layouts in the master
        for layout in master.slide_layouts:
            # Get background fill
            background = layout.background
            
            # Set background color if specified
            bg_color = style_properties.get('backgroundColor') or style_properties.get('background_color')
            if bg_color and bg_color.startswith('#'):
                r, g, b = hex_to_rgb(bg_color)
                background.fill.solid()
                background.fill.fore_color.rgb = RGBColor(r, g, b)
            
            # Apply style to all placeholders in the layout
            for placeholder in layout.placeholders:
                # Process title placeholders
                if placeholder.placeholder_format.type == 1:  # Title placeholder
                    # Set text color
                    text_color = style_properties.get('textColor') or style_properties.get('text_color')
                    if text_color and text_color.startswith('#'):
                        r, g, b = hex_to_rgb(text_color)
                        if placeholder.has_text_frame:
                            for paragraph in placeholder.text_frame.paragraphs:
                                paragraph.font.color.rgb = RGBColor(r, g, b)
                    
                    # Set font family
                    font_family = style_properties.get('fontFamily') or style_properties.get('font_family')
                    if font_family and placeholder.has_text_frame:
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.name = font_family.split(',')[0].strip()
                    
                    # Set title font size
                    title_font_size = style_properties.get('titleFontSize') or style_properties.get('title_font_size')
                    if title_font_size and placeholder.has_text_frame:
                        size_value = title_font_size
                        if isinstance(size_value, str):
                            # Try to extract numeric value from string like "44pt"
                            size_value = size_value.replace('pt', '').replace('px', '')
                            try:
                                size_value = float(size_value)
                            except ValueError:
                                # Fallback to default if can't convert
                                size_value = 44
                        
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.size = Pt(size_value)
                
                # Process body placeholders
                elif placeholder.placeholder_format.type in [2, 3, 7]:  # Body, content or text placeholder
                    # Set text color
                    text_color = style_properties.get('textColor') or style_properties.get('text_color')
                    if text_color and text_color.startswith('#') and placeholder.has_text_frame:
                        r, g, b = hex_to_rgb(text_color)
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.color.rgb = RGBColor(r, g, b)
                    
                    # Set font family
                    font_family = style_properties.get('fontFamily') or style_properties.get('font_family')
                    if font_family and placeholder.has_text_frame:
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.name = font_family.split(',')[0].strip()
                    
                    # Set content font size
                    content_font_size = style_properties.get('contentFontSize') or style_properties.get('content_font_size')
                    if content_font_size and placeholder.has_text_frame:
                        size_value = content_font_size
                        if isinstance(size_value, str):
                            # Try to extract numeric value from string like "28pt"
                            size_value = size_value.replace('pt', '').replace('px', '')
                            try:
                                size_value = float(size_value)
                            except ValueError:
                                # Fallback to default if can't convert
                                size_value = 28
                        
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.size = Pt(size_value)

    return prs

async def create_base_template(template_style, output_path, style_properties=None):
    """Create a basic template in the specified style."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Get template colors based on style
        colors = get_template_colors(template_style)
        
        # Override with custom style properties if provided
        if style_properties:
            primary_color = style_properties.get('primaryColor') or style_properties.get('primary_color')
            if primary_color and primary_color.startswith('#'):
                colors['primary'] = primary_color
            
            secondary_color = style_properties.get('secondaryColor') or style_properties.get('secondary_color')
            if secondary_color and secondary_color.startswith('#'):
                colors['secondary'] = secondary_color
            
            accent_color = style_properties.get('accentColor') or style_properties.get('accent_color')
            if accent_color and accent_color.startswith('#'):
                colors['accent'] = accent_color
            
            bg_color = style_properties.get('backgroundColor') or style_properties.get('background_color')
            if bg_color and bg_color.startswith('#'):
                colors['background'] = bg_color
            
            text_color = style_properties.get('textColor') or style_properties.get('text_color')
            if text_color and text_color.startswith('#'):
                colors['text'] = text_color
        
        # Helper function to convert hex color to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        
        # Create a blank presentation
        prs = Presentation()
        
        # Set slide size to 16:9 aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide layout
        title_slide_layout = prs.slide_layouts[0]
        
        # Apply style to title slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = f"{template_style.capitalize()} Template"
        title_slide.placeholders[1].text = "Created with Slideflow"
        
        # Style title text
        title = title_slide.shapes.title
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.name = colors['fonts']['title']
        title.text_frame.paragraphs[0].font.size = colors['font_sizes']['title']
        title.text_frame.paragraphs[0].font.color.rgb = colors['colors']['title']
        title.text_frame.paragraphs[0].font.bold = colors['styles']['title_bold']
        
        # Style subtitle text
        subtitle = title_slide.placeholders[1]
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle.text_frame.paragraphs[0].font.name = colors['fonts']['body']
        subtitle.text_frame.paragraphs[0].font.size = colors['font_sizes']['body']
        subtitle.text_frame.paragraphs[0].font.color.rgb = colors['colors']['body']
        subtitle.text_frame.paragraphs[0].font.bold = colors['styles']['body_bold']
        
        # Create content slide layout
        content_slide_layout = prs.slide_layouts[1]
        
        # Apply style to content slide
        content_slide = prs.slides.add_slide(content_slide_layout)
        content_slide.shapes.title.text = "Content Slide Example"
        content_text = content_slide.placeholders[1]
        
        # Add sample bullet points
        content_text.text = "Sample Content"
        
        # Style title text
        title = content_slide.shapes.title
        title.text_frame.paragraphs[0].font.name = colors['fonts']['title']
        title.text_frame.paragraphs[0].font.size = colors['font_sizes']['title']
        title.text_frame.paragraphs[0].font.color.rgb = colors['colors']['title']
        title.text_frame.paragraphs[0].font.bold = colors['styles']['title_bold']
        
        # Style content text
        content_text.text_frame.paragraphs[0].font.name = colors['fonts']['body']
        content_text.text_frame.paragraphs[0].font.size = colors['font_sizes']['body']
        content_text.text_frame.paragraphs[0].font.color.rgb = colors['colors']['body']
        content_text.text_frame.paragraphs[0].font.bold = colors['styles']['body_bold']
        
        # Save the template
        prs.save(output_path)
        
        return output_path
    except Exception as e:
        print(f"Error creating base template: {str(e)}")
        raise

async def export_to_docx(content_data: Dict[str, Any], output_path: str) -> str:
    """
    Export presentation content to a Word document.
    
    Args:
        content_data: The presentation content data
        output_path: Path to save the document
        
    Returns:
        Path to the created document
    """
    try:
        from docx import Document
        
        # Create a new Word document
        doc = Document()
        
        # Add title
        if "title" in content_data:
            doc.add_heading(content_data["title"], level=0)
        
        # Add description if present
        if "description" in content_data:
            doc.add_paragraph(content_data["description"])
        
        # Add content from slides
        if "slides" in content_data:
            for slide_data in content_data["slides"]:
                # Make a copy of slide_data to avoid modifying the original
                processed_slide_data = slide_data.copy()
                
                # Remove image_prompt as we don't want to include it in the document
                if "image_prompt" in processed_slide_data:
                    processed_slide_data.pop("image_prompt")
                
                # Skip title slide as we've already added it
                if processed_slide_data.get("type") == "title" and processed_slide_data.get("title") == content_data.get("title"):
                    continue
                
                # Add slide title as heading
                if "title" in processed_slide_data:
                    doc.add_heading(processed_slide_data["title"], level=1)
                
                # Add content based on slide type
                if processed_slide_data.get("type") == "quote":
                    if "quote" in processed_slide_data:
                        quote_para = doc.add_paragraph(f'"{processed_slide_data["quote"]}"', style='Intense Quote')
                        if "author" in processed_slide_data:
                            doc.add_paragraph(f"— {processed_slide_data['author']}")
                
                elif "bullets" in processed_slide_data:
                    for bullet in processed_slide_data["bullets"]:
                        doc.add_paragraph(bullet, style='List Bullet')
                
                elif "content" in processed_slide_data:
                    doc.add_paragraph(processed_slide_data["content"])
                
                elif "left_content" in processed_slide_data or "right_content" in processed_slide_data:
                    if "left_content" in processed_slide_data:
                        doc.add_paragraph("Left Column:", style='Strong')
                        doc.add_paragraph(processed_slide_data["left_content"])
                    
                    if "right_content" in processed_slide_data:
                        doc.add_paragraph("Right Column:", style='Strong')
                        doc.add_paragraph(processed_slide_data["right_content"])
                
                # Add notes if present
                if "notes" in processed_slide_data:
                    doc.add_paragraph(processed_slide_data["notes"], style='Quote')
                
                # Add a page break between slides
                doc.add_page_break()
        
        # Save the document
        doc.save(output_path)
        return output_path
        
    except ImportError:
        print("python-docx library not installed. Install it with: pip install python-docx")
        raise
    except Exception as e:
        print(f"Error creating Word document: {str(e)}")
        raise

async def export_to_pdf(content_data: Dict[str, Any], template_style: str, output_path: str, pptx_path: Optional[str] = None) -> str:
    """
    Export presentation content to a PDF document.
    
    Args:
        content_data: The presentation content data
        template_style: The template style to use
        output_path: Path to save the PDF
        pptx_path: Optional path to an existing PPTX file to convert (if None, one will be created)
        
    Returns:
        Path to the created PDF
    """
    try:
        # First create a PowerPoint presentation if not provided
        generated_temp_file = False
        if not pptx_path:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
                pptx_path = tmp_pptx.name
            
            # Create a deep copy of content_data to ensure we don't modify the original
            import copy
            modified_content_data = copy.deepcopy(content_data)
            
            # Remove image prompts from all slides
            if "slides" in modified_content_data:
                for slide in modified_content_data["slides"]:
                    if "image_prompt" in slide:
                        slide.pop("image_prompt")
            
            # Generate the presentation (images will still be included)
            await create_presentation(modified_content_data, template_style, pptx_path)
            generated_temp_file = True
        
        # Convert to PDF using LibreOffice (server must have LibreOffice installed)
        # For Windows, adjust command if needed
        if os.name == 'nt':  # Windows
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path):
                    libreoffice_path = f'"{path}"'
                    break
            
            if not libreoffice_path:
                raise Exception("LibreOffice not found. Please install it to convert to PDF.")
            
            cmd = f'{libreoffice_path} --headless --convert-to pdf --outdir "{os.path.dirname(output_path)}" "{pptx_path}"'
        else:  # Linux/Mac
            cmd = f'soffice --headless --convert-to pdf --outdir "{os.path.dirname(output_path)}" "{pptx_path}"'
        
        # Run the conversion command
        process = await asyncio.create_subprocess_shell(
            cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE
        )
        
        stdout, stderr = await process.communicate()
        
        if process.returncode != 0:
            print(f"PDF conversion error: {stderr.decode()}")
            raise Exception(f"Failed to convert to PDF: {stderr.decode()}")
        
        # Get the output PDF path
        pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'
        
        # Move to desired output location if different
        if pdf_path != output_path:
            import shutil
            shutil.move(pdf_path, output_path)
        
        # Clean up temporary file if we created it
        if generated_temp_file:
            os.unlink(pptx_path)
        
        return output_path
        
    except Exception as e:
        print(f"Error creating PDF: {str(e)}")
        raise

def get_template_colors(template_style):
    """Get colors and style information for the specified template style."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    # Helper function to convert hex to RGB
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    # Define template styles
    templates = {
        "corporate": {
            "primary": "#0F62FE",
            "secondary": "#6F6F6F",
            "accent": "#4589FF",
            "background": "#FFFFFF",
            "text": "#002B5B",
            "fonts": {
                "title": "Calibri",
                "body": "Arial"
            },
            "font_sizes": {
                "title": Pt(40),
                "body": Pt(24)
            },
            "styles": {
                "title_bold": True,
                "body_bold": False
            }
        },
        "creative": {
            "primary": "#FF3366",
            "secondary": "#9C27B0",
            "accent": "#FFCC00",
            "background": "#F5F5F5",
            "text": "#5A189A",
            "fonts": {
                "title": "Poppins",
                "body": "Open Sans"
            },
            "font_sizes": {
                "title": Pt(44),
                "body": Pt(26)
            },
            "styles": {
                "title_bold": True,
                "body_bold": False
            }
        },
        "minimalist": {
            "primary": "#212121",
            "secondary": "#757575",
            "accent": "#BDBDBD",
            "background": "#FFFFFF",
            "text": "#444444",
            "fonts": {
                "title": "Helvetica",
                "body": "Arial"
            },
            "font_sizes": {
                "title": Pt(40),
                "body": Pt(24)
            },
            "styles": {
                "title_bold": False,
                "body_bold": False
            }
        },
        "academic": {
            "primary": "#006064",
            "secondary": "#00897B",
            "accent": "#4DD0E1",
            "background": "#FAF3E0",
            "text": "#1C3D6E",
            "fonts": {
                "title": "Times New Roman",
                "body": "Georgia"
            },
            "font_sizes": {
                "title": Pt(38),
                "body": Pt(22)
            },
            "styles": {
                "title_bold": True,
                "body_bold": False
            }
        },
        "marketing": {
            "primary": "#FF5722",
            "secondary": "#FF9800",
            "accent": "#FFC107",
            "background": "#FFFFFF",
            "text": "#E63946",
            "fonts": {
                "title": "Impact",
                "body": "Arial"
            },
            "font_sizes": {
                "title": Pt(44),
                "body": Pt(26)
            },
            "styles": {
                "title_bold": True,
                "body_bold": False
            }
        },
        "techStartup": {
            "primary": "#00A8E8",
            "secondary": "#1E1E1E",
            "accent": "#2ECC71",
            "background": "#1E1E1E",
            "text": "#00A8E8",
            "fonts": {
                "title": "Arial",
                "body": "Verdana"
            },
            "font_sizes": {
                "title": Pt(42),
                "body": Pt(24)
            },
            "styles": {
                "title_bold": True,
                "body_bold": False
            }
        }
    }
    
    # Use corporate as default if template not found
    template_data = templates.get(template_style, templates["corporate"])
    
    # Convert hex colors to RGB
    colors = {
        "primary": template_data["primary"],
        "secondary": template_data["secondary"],
        "accent": template_data["accent"],
        "background": template_data["background"],
        "text": template_data["text"],
        "fonts": template_data["fonts"],
        "font_sizes": template_data["font_sizes"],
        "styles": template_data["styles"],
        "colors": {
            "title": RGBColor(*hex_to_rgb(template_data["text"])),
            "body": RGBColor(*hex_to_rgb(template_data["text"])),
            "background": RGBColor(*hex_to_rgb(template_data["background"])),
            "accent": RGBColor(*hex_to_rgb(template_data["accent"]))
        }
    }
    
    return colors
        
async def create_base_template(template_style, output_path, style_properties=None):
    """Create a basic template in the specified style."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Get template colors based on style
        colors = get_template_colors(template_style)
        
        # Override with custom style properties if provided
        if style_properties:
            primary_color = style_properties.get('primaryColor') or style_properties.get('primary_color')
            if primary_color and primary_color.startswith('#'):
                colors['primary'] = primary_color
            
            secondary_color = style_properties.get('secondaryColor') or style_properties.get('secondary_color')
            if secondary_color and secondary_color.startswith('#'):
                colors['secondary'] = secondary_color
            
            accent_color = style_properties.get('accentColor') or style_properties.get('accent_color')
            if accent_color and accent_color.startswith('#'):
                colors['accent'] = accent_color
            
            bg_color = style_properties.get('backgroundColor') or style_properties.get('background_color')
            if bg_color and bg_color.startswith('#'):
                colors['background'] = bg_color
            
            text_color = style_properties.get('textColor') or style_properties.get('text_color')
            if text_color and text_color.startswith('#'):
                colors['text'] = text_color
        
        # Helper function to convert hex color to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        
        # Create a blank presentation
        prs = Presentation()
        
        # Set slide size to 16:9 aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide layout
        title_slide_layout = prs.slide_layouts[0]
        
        # Apply style to title slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = f"{template_style.capitalize()} Template"
        title_slide.placeholders[1].text = "Created with Slideflow"
        
        # Style title text
        title = title_slide.shapes.title
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.name = colors['fonts']['title']
        title.text_frame.paragraphs[0].font.size = colors['font_sizes']['title']
        title.text_frame.paragraphs[0].font.color.rgb = colors['colors']['title']
        title.text_frame.paragraphs[0].font.bold = colors['styles']['title_bold']
        
        # Style subtitle text
        subtitle = title_slide.placeholders[1]
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle.text_frame.paragraphs[0].font.name = colors['fonts']['body']
        subtitle.text_frame.paragraphs[0].font.size = colors['font_sizes']['body']
        subtitle.text_frame.paragraphs[0].font.color.rgb = colors['colors']['body']
        subtitle.text_frame.paragraphs[0].font.bold = colors['styles']['body_bold']
        
        # Create content slide layout
        content_slide_layout = prs.slide_layouts[1]
        
        # Apply style to content slide
        content_slide = prs.slides.add_slide(content_slide_layout)
        content_slide.shapes.title.text = "Content Slide Example"
        content_text = content_slide.placeholders[1]
        
        # Add sample bullet points
        content_text.text = "Sample Content"
        
        # Style title text
        title = content_slide.shapes.title
        title.text_frame.paragraphs[0].font.name = colors['fonts']['title']
        title.text_frame.paragraphs[0].font.size = colors['font_sizes']['title']
        title.text_frame.paragraphs[0].font.color.rgb = colors['colors']['title']
        title.text_frame.paragraphs[0].font.bold = colors['styles']['title_bold']
        
        # Style content text
        content_text.text_frame.paragraphs[0].font.name = colors['fonts']['body']
        content_text.text_frame.paragraphs[0].font.size = colors['font_sizes']['body']
        content_text.text_frame.paragraphs[0].font.color.rgb = colors['colors']['body']
        content_text.text_frame.paragraphs[0].font.bold = colors['styles']['body_bold']
        
        # Save the template
        prs.save(output_path)
        
        return output_path
    except Exception as e:
        print(f"Error creating base template: {str(e)}")
        raise 